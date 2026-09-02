import os
import csv
from pptx import Presentation


def extract_text_from_slide(slide):
    title_text = slide.shapes.title.text if slide.shapes.title else ""

    text_box_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            text_box_text += shape.text

    return title_text, text_box_text


def process_pptx_file(file_path, output_csv_path):
    presentation = Presentation(file_path)

    header = ['videoName', 'Cx1', 'Cy1', 'Cx2', 'Cy2', 'Cx3', 'Cy3', 'Cx4', 'Cy4']
    rows = [header]

    for i, slide in enumerate(presentation.slides):
        title_text, text_box_text = extract_text_from_slide(slide)
        title_text = title_text + '.mp4'

        row = [title_text, text_box_text]
        rows.append(row)

    with open(output_csv_path, 'w', newline='', encoding='utf-8') as csvfile:
        csv_writer = csv.writer(csvfile)
        csv_writer.writerows(rows)

if __name__ == "__main__":
    pptx_file_path = r'E:\videos\cologne_yijiang\extracted_snaps\extracted_snaps\all_snapshots_20231120_034947.pptx'
    output_csv_file_path = r'E:\videos\cologne_yijiang\extracted_snaps\extracted_snaps\output.csv'

    process_pptx_file(pptx_file_path, output_csv_file_path)
