import os
import cv2
from pptx import Presentation
from pptx.util import Inches
import tempfile
import datetime

timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")


def get_single_frame(video_path):
    cap = cv2.VideoCapture(video_path)
    ret, frame = cap.read()
    cap.release()
    return frame if ret else None


def process_videos(input_folder, output_folder):
    video_files = [f for f in os.listdir(input_folder) if f.endswith(('.avi', '.mp4'))]

    presentation = Presentation()

    for video_file in video_files:
        video_path = os.path.join(input_folder, video_file)
        frame = get_single_frame(video_path)

        if frame is not None:
            # Add a slide for the snapshot in the PowerPoint file
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Use blank layout

            # Add title with video file name
            title = slide.shapes.title
            title.text = os.path.splitext(video_file)[0]

            # Save the image to a temporary file
            temp_image_path = os.path.join(tempfile.gettempdir(), "temp_snapshot.png")
            cv2.imwrite(temp_image_path, frame)

            # Add image to the slide
            slide.shapes.add_picture(temp_image_path, Inches(0), Inches(2), width=Inches(10))

            # Delete the temporary file
            os.remove(temp_image_path)

    # Save PowerPoint file
    ppt_output_path = os.path.join(output_folder, f"all_snapshots_{timestamp}.pptx")
    presentation.save(ppt_output_path)
    print(f"PowerPoint file saved: {ppt_output_path}")


if __name__ == "__main__":
    input_folder = r'E:\videos\7'
    output_folder = os.path.join(input_folder, 'extracted_snaps')

    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    process_videos(input_folder, output_folder)
